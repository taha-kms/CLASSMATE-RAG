"""
PPTX loader for CLASSMATE-RAG.

Extracts readable text from slides (titles, shapes, tables) and speaker notes.
Returns a list of (page_index, text) tuples, where page_index is the 1-based slide number.

Notes:
- This loader focuses on textual content. Images are ignored (OCR handled elsewhere if needed).
- Tables are read cell-by-cell in row order.
- Speaker notes (if present) are appended after slide text, prefixed for clarity.
"""

from __future__ import annotations

from pathlib import Path
from typing import List, Tuple

from pptx import Presentation


def _shape_text(shape) -> str:
    """Extract text from a shape if it has a text_frame."""
    try:
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            return shape.text or ""
    except Exception:
        pass
    return ""


def _table_text(shape) -> str:
    """Extract text from a table, row by row."""
    try:
        if hasattr(shape, "has_table") and shape.has_table:
            table = shape.table
            rows = []
            for r in table.rows:
                cells = []
                for c in r.cells:
                    cells.append((c.text or "").strip())
                rows.append(" | ".join(filter(None, cells)))
            return "\n".join(filter(None, rows))
    except Exception:
        pass
    return ""


def _notes_text(slide) -> str:
    """Extract speaker notes if present."""
    try:
        notes = slide.notes_slide
        if notes and notes.notes_text_frame:
            txt = notes.notes_text_frame.text or ""
            txt = txt.strip()
            if txt:
                return f"[Notes]\n{txt}"
    except Exception:
        # No notes or not available
        pass
    return ""


def _collect_slide_text(slide) -> str:
    blocks: List[str] = []

    # Title placeholder if available
    try:
        if slide.shapes.title and slide.shapes.title.text:
            blocks.append(slide.shapes.title.text.strip())
    except Exception:
        pass

    # Other shapes (text boxes etc.)
    for shp in slide.shapes:
        t = _shape_text(shp).strip()
        if t:
            blocks.append(t)
        # Tables
        tt = _table_text(shp).strip()
        if tt:
            blocks.append(tt)

    # Speaker notes at the end (clearly separated)
    nt = _notes_text(slide).strip()
    if nt:
        blocks.append(nt)

    # Deduplicate consecutive identical blocks and join
    cleaned: List[str] = []
    last = None
    for b in blocks:
        if b and b != last:
            cleaned.append(b)
        last = b
    return "\n\n".join(cleaned)


def load_pptx_slides(path: str | Path) -> List[Tuple[int, str]]:
    """
    Read a .pptx file and return a list of (page_index, text) tuples.
    page_index is 1-based to match human-friendly slide numbering.

    Example:
        slides = load_pptx_slides("lecture01.pptx")
        for page, text in slides:
            print(page, text[:200])
    """
    p = Path(path)
    if not p.exists() or p.suffix.lower() not in (".pptx",):
        raise ValueError(f"Expected an existing .pptx file, got: {path}")

    pres = Presentation(str(p))
    out: List[Tuple[int, str]] = []

    for i, slide in enumerate(pres.slides, start=1):
        txt = _collect_slide_text(slide).strip()
        if not txt:
            # Empty slide — keep a placeholder to preserve ordering
            out.append((i, ""))  # caller may skip empty chunks later
        else:
            out.append((i, txt))

    return out
