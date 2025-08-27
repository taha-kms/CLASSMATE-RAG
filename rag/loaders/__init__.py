"""
Unified loader interface.

Public:
    infer_doc_type_from_path(path) -> str
    load_document_by_type(path, doc_type, enable_ocr=False) -> list[(page, text)]

This module implements loaders for:
    - txt, md
    - html (readability)
    - csv (bullets)
    - epub
    - pdf, docx, pptx (basic, dependency-backed)

Notes
- All loaders return a list of (page_number:int, text:str).
- For multi-part formats (pptx slides, epub chapters), page numbers increase per part.
"""

from __future__ import annotations

import os
from pathlib import Path
from typing import List, Tuple

from .html_readable import load_html_readable
from .csv_bullets import load_csv_bullets
from .epub_loader import load_epub

# Optional deps for document types
try:
    import pypdf  # type: ignore
except Exception:  # pragma: no cover
    pypdf = None  # type: ignore

try:
    import docx  # python-docx  # type: ignore
except Exception:  # pragma: no cover
    docx = None  # type: ignore

try:
    import pptx  # python-pptx  # type: ignore
except Exception:  # pragma: no cover
    pptx = None  # type: ignore


# -----------------------
# Type inference
# -----------------------

def infer_doc_type_from_path(path: str | Path) -> str:
    ext = Path(path).suffix.lower().lstrip(".")
    if ext in {"htm", "html"}:
        return "html"
    if ext == "csv":
        return "csv"
    if ext == "epub":
        return "epub"
    if ext in {"md", "markdown"}:
        return "md"
    if ext in {"txt", "text"}:
        return "txt"
    if ext == "pdf":
        return "pdf"
    if ext in {"docx"}:
        return "docx"
    if ext in {"pptx", "ppt"}:
        return "pptx"
    return "other"


# -----------------------
# Primitive loaders
# -----------------------

def _load_txt(path: Path) -> List[Tuple[int, str]]:
    txt = path.read_text(encoding="utf-8", errors="ignore")
    txt = (txt or "").strip()
    return [(1, txt)] if txt else []


def _load_md(path: Path) -> List[Tuple[int, str]]:
    # Keep markdown markup; later chunker/tokenizer deals with structure.
    md = path.read_text(encoding="utf-8", errors="ignore")
    md = (md or "").strip()
    return [(1, md)] if md else []


def _load_pdf(path: Path) -> List[Tuple[int, str]]:
    if pypdf is None:
        raise RuntimeError("pypdf is not installed; cannot load PDF")
    reader = pypdf.PdfReader(str(path))  # type: ignore
    pages: List[Tuple[int, str]] = []
    for i, page in enumerate(reader.pages, start=1):  # type: ignore
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        text = text.strip()
        if text:
            pages.append((i, text))
    return pages


def _load_docx(path: Path) -> List[Tuple[int, str]]:
    if docx is None:
        raise RuntimeError("python-docx is not installed; cannot load DOCX")
    doc = docx.Document(str(path))  # type: ignore
    paras = [p.text.strip() for p in doc.paragraphs]  # type: ignore
    text = "\n".join([p for p in paras if p])
    text = text.strip()
    return [(1, text)] if text else []


def _load_pptx(path: Path) -> List[Tuple[int, str]]:
    if pptx is None:
        raise RuntimeError("python-pptx is not installed; cannot load PPTX")
    prs = pptx.Presentation(str(path))  # type: ignore
    pages: List[Tuple[int, str]] = []
    for i, slide in enumerate(prs.slides, start=1):  # type: ignore
        chunks: List[str] = []
        for shape in slide.shapes:  # type: ignore
            try:
                if hasattr(shape, "text"):
                    t = (shape.text or "").strip()
                    if t:
                        chunks.append(t)
            except Exception:
                continue
        txt = "\n".join(chunks).strip()
        if txt:
            pages.append((i, txt))
    return pages


# -----------------------
# Unified entrypoint
# -----------------------

def load_document_by_type(
    path: str | Path,
    doc_type: str,
    *,
    enable_ocr: bool = False,  # kept for API compatibility; not used here
) -> List[Tuple[int, str]]:
    """
    Route to the specific loader based on doc_type (string).
    """
    p = Path(path).expanduser().resolve()
    t = (doc_type or infer_doc_type_from_path(p)).lower()

    if t == "txt":
        return _load_txt(p)
    if t == "md":
        return _load_md(p)
    if t == "html":
        return load_html_readable(p)
    if t == "csv":
        return load_csv_bullets(p)
    if t == "epub":
        return load_epub(p)
    if t == "pdf":
        return _load_pdf(p)
    if t == "docx":
        return _load_docx(p)
    if t == "pptx":
        return _load_pptx(p)

    # Fallback: try plain text
    return _load_txt(p)
